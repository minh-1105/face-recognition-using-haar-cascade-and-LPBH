from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(243, 241, 235)
BG_SOFT = RGBColor(235, 230, 221)
TEXT = RGBColor(18, 18, 18)
MUTED = RGBColor(93, 90, 84)
ACCENT = RGBColor(217, 71, 43)
LINE = RGBColor(212, 204, 194)
WHITE = RGBColor(255, 255, 255)

DISPLAY_FONT = "Aptos Display"
BODY_FONT = "Aptos"


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.18), Inches(0.18), Inches(12.97), Inches(7.14)
    )
    frame.fill.background()
    frame.line.color.rgb = LINE
    frame.line.width = Pt(1)

    halo = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.45), Inches(-0.45), Inches(3.3), Inches(3.3)
    )
    halo.fill.background()
    halo.line.color.rgb = LINE
    halo.line.width = Pt(1)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = BG_SOFT
    panel.fill.transparency = 0.83
    panel.line.fill.background()
    slide.shapes._spTree.remove(panel._element)
    slide.shapes._spTree.insert(2, panel._element)


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text="",
    size=18,
    color=TEXT,
    bold=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
    upper=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text.upper() if upper else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for idx, item in enumerate(items):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(10)
        paragraph.bullet = True
        for run in paragraph.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(size)
            run.font.color.rgb = MUTED

    return box


def add_card(slide, x, y, w, h, title, body):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.fill.transparency = 0.62
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)

    add_textbox(slide, x + 0.18, y + 0.15, w - 0.36, 0.28, title, size=11, color=ACCENT, bold=True, upper=True)
    add_textbox(slide, x + 0.18, y + 0.44, w - 0.36, h - 0.55, body, size=16, color=MUTED)
    return shape


def add_rail(slide, number, label):
    add_textbox(slide, 11.8, 6.5, 0.9, 0.25, f"{number:02d}", size=11, color=ACCENT, bold=True, upper=True)
    add_textbox(slide, 11.15, 6.76, 1.6, 0.25, label, size=10, color=MUTED, upper=True, align=PP_ALIGN.RIGHT)


def add_title_block(slide, eyebrow, title, subtitle=None):
    add_textbox(slide, 0.72, 0.6, 3.6, 0.25, eyebrow, size=11, color=ACCENT, bold=True, upper=True)
    add_textbox(slide, 0.72, 1.0, 6.6, 1.1, title, size=28, color=TEXT, bold=True, font=DISPLAY_FONT)
    if subtitle:
        add_textbox(slide, 0.72, 2.0, 6.3, 0.65, subtitle, size=15, color=MUTED)


def add_chip(slide, x, y, text):
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(1.45), Inches(0.42)
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = WHITE
    chip.fill.transparency = 0.55
    chip.line.color.rgb = LINE
    chip.line.width = Pt(1)
    add_textbox(slide, x + 0.12, y + 0.08, 1.2, 0.2, text, size=11, color=MUTED)


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = LINE
    line.line.width = Pt(1.25)
    line.line.end_arrowhead = True
    return line


def slide_intro():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.72, 1.0, 3.2, 0.25, "Assignment Group · Computer Vision", size=11, color=ACCENT, bold=True, upper=True)
    add_textbox(slide, 0.72, 1.34, 3.0, 0.22, "Pipeline phát hiện khuôn mặt", size=11, color=MUTED, upper=True)
    add_textbox(slide, 0.72, 1.85, 6.9, 1.7, "Haar Cascade\ntrong ứng dụng điểm danh", size=31, color=TEXT, bold=True, font=DISPLAY_FONT)
    add_textbox(
        slide,
        0.72,
        3.45,
        6.0,
        0.7,
        "Trọng tâm: cơ chế đặc trưng Haar, integral image, AdaBoost, cascade classifier và vị trí của thuật toán trong ứng dụng.",
        size=16,
        color=MUTED,
    )
    add_chip(slide, 0.72, 4.35, "8 slides")
    add_chip(slide, 2.3, 4.35, "Kỹ thuật")
    add_chip(slide, 3.88, 4.35, "Trọng tâm")
    add_rail(slide, 1, "INTRO")


def slide_concept():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_block(
        slide,
        "Haar Cascade",
        "Bài toán mà Haar Cascade giải quyết",
        "Bộ phân loại phải quyết định mỗi cửa sổ ảnh có chứa khuôn mặt hay không.",
    )
    add_bullets(
        slide,
        0.72,
        2.8,
        5.5,
        2.4,
        [
            "Ảnh đầu vào thường được chuyển sang grayscale để giảm số kênh cần xử lý.",
            "Thuật toán trượt một cửa sổ qua nhiều vị trí và nhiều tỉ lệ kích thước.",
            "Tại từng cửa sổ, hệ thống tính tập feature Haar rồi phân loại mặt / không mặt.",
            "Đầu ra là các bounding box dùng cho bước nhận diện hoặc theo dõi tiếp theo.",
        ],
        size=17,
    )
    add_card(
        slide,
        7.0,
        1.65,
        5.0,
        1.4,
        "Giả định thị giác",
        "Khuôn mặt chính diện có mẫu tương phản khá ổn định: mắt tối hơn má, sống mũi sáng hơn hai bên, nên có thể mô tả bằng các vùng chữ nhật sáng - tối.",
    )

    positions = [(7.0, 3.45), (8.73, 3.45), (10.46, 3.45)]
    for idx, (x, y) in enumerate(positions):
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(1.45), Inches(1.45))
        frame.fill.background()
        frame.line.color.rgb = LINE
        frame.line.width = Pt(1)
        if idx == 0:
            dark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.72), Inches(1.45))
            light = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.72), Inches(y), Inches(0.73), Inches(1.45))
            dark.fill.solid()
            dark.fill.fore_color.rgb = TEXT
            dark.line.fill.background()
            light.fill.solid()
            light.fill.fore_color.rgb = BG
            light.line.fill.background()
        elif idx == 1:
            dark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(1.45), Inches(0.72))
            light = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y + 0.72), Inches(1.45), Inches(0.73))
            dark.fill.solid()
            dark.fill.fore_color.rgb = TEXT
            dark.line.fill.background()
            light.fill.solid()
            light.fill.fore_color.rgb = BG
            light.line.fill.background()
        else:
            part_w = 1.45 / 3
            for j, color in enumerate((BG, TEXT, BG)):
                part = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + j * part_w), Inches(y), Inches(part_w), Inches(1.45)
                )
                part.fill.solid()
                part.fill.fore_color.rgb = color
                part.line.fill.background()

    add_textbox(slide, 7.0, 5.2, 5.0, 0.45, "Ba mẫu cơ bản: cạnh dọc, cạnh ngang và dải sáng - tối - sáng.", size=14, color=MUTED)
    add_rail(slide, 2, "CONCEPT")


def slide_features():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_block(slide, "Core Mechanics", "Đặc trưng Haar và integral image")
    add_card(
        slide,
        0.72,
        1.95,
        3.9,
        1.5,
        "Haar-like feature",
        "Mỗi feature là hiệu giữa tổng mức xám của các vùng chữ nhật. Giá trị này phản ánh sự thay đổi cường độ sáng theo một mẫu hình học cố định.",
    )
    add_card(
        slide,
        4.82,
        1.95,
        3.9,
        1.5,
        "Integral image",
        "Integral image lưu tổng tích lũy đến mỗi tọa độ. Nhờ đó, tổng của một hình chữ nhật chỉ cần 4 giá trị biên A, B, C, D.",
    )

    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(3.75), Inches(8.0), Inches(0.52))
    strip.fill.background()
    strip.line.color.rgb = LINE
    strip.line.width = Pt(1)
    add_textbox(slide, 0.95, 3.89, 2.8, 0.2, "sum(rect) = D - B - C + A", size=16, color=TEXT, bold=True, font=DISPLAY_FONT)
    add_textbox(
        slide,
        4.0,
        3.89,
        4.4,
        0.2,
        "Không cần cộng từng pixel cho từng feature",
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )

    demo = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.05), Inches(1.95), Inches(3.0), Inches(2.95))
    demo.fill.background()
    demo.line.color.rgb = LINE
    demo.line.width = Pt(1)

    for i in range(1, 5):
        h_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.05), Inches(1.95 + i * 0.58), Inches(12.05), Inches(1.95 + i * 0.58))
        h_line.line.color.rgb = LINE
        h_line.line.width = Pt(1)
        v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.05 + i * 0.6), Inches(1.95), Inches(9.05 + i * 0.6), Inches(4.9))
        v_line.line.color.rgb = LINE
        v_line.line.width = Pt(1)

    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.6), Inches(2.45), Inches(1.9), Inches(1.55))
    rect.fill.solid()
    rect.fill.fore_color.rgb = ACCENT
    rect.fill.transparency = 0.82
    rect.line.color.rgb = ACCENT
    rect.line.width = Pt(1.5)

    for label, x, y in (("A", 9.35, 2.2), ("B", 11.35, 2.2), ("C", 9.35, 4.05), ("D", 11.35, 4.05)):
        add_textbox(slide, x, y, 0.25, 0.2, label, size=14, color=TEXT, bold=True, font=DISPLAY_FONT)

    add_textbox(
        slide,
        0.72,
        4.75,
        7.6,
        0.6,
        "Ý nghĩa thực tế: một frame có thể chứa hàng nghìn cửa sổ và mỗi cửa sổ có rất nhiều feature. Integral image giúp việc tính toán đủ nhanh để dùng realtime.",
        size=15,
        color=MUTED,
    )
    add_rail(slide, 3, "FEATURES")


def slide_training():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_block(slide, "Training", "AdaBoost chọn feature và ngưỡng")

    steps = [
        ("01", "Mẫu dương / âm", "Tập huấn luyện gồm ảnh có mặt và ảnh không có mặt ở cùng kích thước chuẩn."),
        ("02", "Sinh nhiều feature", "Từ mỗi cửa sổ chuẩn, hệ thống sinh ra số lượng rất lớn feature hình chữ nhật."),
        ("03", "AdaBoost tối ưu", "AdaBoost chọn các feature có sai số thấp và tăng trọng số cho các mẫu đang bị phân loại sai."),
        ("04", "Weak classifier", "Mỗi feature được gắn với ngưỡng và chiều so sánh để tạo một bộ phân loại yếu."),
    ]

    start_x = 0.72
    for idx, (num, title, body) in enumerate(steps):
        x = start_x + idx * 3.0
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(2.45), Inches(x + 2.4), Inches(2.45))
        line.line.color.rgb = LINE
        line.line.width = Pt(1)
        add_textbox(slide, x, 2.05, 0.4, 0.2, num, size=11, color=ACCENT, bold=True, upper=True)
        add_textbox(slide, x, 2.6, 2.4, 0.5, title, size=16, color=TEXT, bold=True, font=DISPLAY_FONT)
        add_textbox(slide, x, 3.1, 2.45, 1.05, body, size=14, color=MUTED)

    add_card(
        slide,
        6.95,
        5.55,
        4.95,
        0.95,
        "Kết quả",
        "Mô hình cuối cùng chỉ giữ một tập nhỏ feature có sức phân biệt cao, thay vì dùng toàn bộ feature đã sinh ra ban đầu.",
    )
    add_rail(slide, 4, "TRAINING")


def slide_cascade():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_block(slide, "Detection Pipeline", "Cascade classifier khi suy luận")
    add_bullets(
        slide,
        0.72,
        2.55,
        5.3,
        2.1,
        [
            "Mỗi cửa sổ ảnh phải vượt qua nhiều stage nối tiếp.",
            "Stage đầu dùng ít feature để loại phần lớn negative window.",
            "Chỉ các cửa sổ còn nghi ngờ mới được chuyển sang stage tiếp theo.",
            "Cửa sổ bị loại ở bất kỳ stage nào sẽ dừng xử lý ngay.",
        ],
        size=17,
    )

    stages = [
        ("Stage 1", "few features", False),
        ("Stage 2", "more filters", False),
        ("Stage 3", "strict check", False),
        ("Face", "accepted window", True),
    ]

    y = 2.2
    for title, note, final in stages:
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.0), Inches(y), Inches(4.8), Inches(0.58))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT if final else WHITE
        bar.fill.transparency = 0.88 if final else 0.68
        bar.line.color.rgb = ACCENT if final else LINE
        bar.line.width = Pt(1)
        add_textbox(slide, 7.25, y + 0.14, 1.8, 0.2, title, size=15, color=TEXT, bold=True, font=DISPLAY_FONT)
        add_textbox(slide, 9.8, y + 0.14, 1.6, 0.2, note, size=12, color=MUTED, align=PP_ALIGN.RIGHT)
        y += 0.75

    add_textbox(slide, 7.0, 5.45, 4.8, 0.22, "Ưu điểm: tốc độ cao vì loại negative window rất sớm.", size=15, color=MUTED)
    add_textbox(slide, 7.0, 5.82, 4.8, 0.22, "Hạn chế: kém ổn định khi mặt nghiêng mạnh, ánh sáng xấu hoặc bị che khuất.", size=15, color=MUTED)
    add_rail(slide, 5, "CASCADE")


def slide_pipeline():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_block(slide, "Application Overview", "Vị trí của Haar Cascade trong ứng dụng")

    nodes = [
        ("RTSP /\nWebcam", 0.72, False),
        ("Grayscale +\nEqualize + Blur", 3.08, False),
        ("Haar Cascade\nphát hiện mặt", 5.44, True),
        ("LBPH\nnhận diện", 8.1, False),
        ("CSV\nđiểm danh", 10.46, False),
    ]

    for idx, (label, x, accent) in enumerate(nodes):
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(3.0), Inches(1.95), Inches(1.0))
        node.fill.solid()
        node.fill.fore_color.rgb = ACCENT if accent else WHITE
        node.fill.transparency = 0.88 if accent else 0.65
        node.line.color.rgb = ACCENT if accent else LINE
        node.line.width = Pt(1)
        add_textbox(
            slide,
            x + 0.12,
            3.24,
            1.7,
            0.45,
            label,
            size=14,
            color=TEXT,
            bold=accent,
            font=DISPLAY_FONT if accent else BODY_FONT,
            align=PP_ALIGN.CENTER,
        )
        if idx < len(nodes) - 1:
            add_arrow(slide, x + 1.95, 3.5, x + 2.25, 3.5)

    add_textbox(
        slide,
        0.72,
        4.95,
        8.1,
        0.55,
        "Trong code hiện tại, ảnh được chuyển sang grayscale, cân bằng histogram, làm mờ Gaussian rồi đưa vào detectMultiScale(scaleFactor=1.2, minNeighbors=5, minSize=(90, 90)).",
        size=15,
        color=MUTED,
    )
    add_textbox(
        slide,
        0.72,
        5.48,
        7.8,
        0.35,
        "Haar Cascade chỉ đảm nhiệm phát hiện vùng mặt; phần nhận diện danh tính được thực hiện riêng bằng LBPH.",
        size=15,
        color=MUTED,
    )
    add_rail(slide, 6, "PIPELINE")


def slide_usage():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_block(slide, "App Usage", "Chức năng chính của ứng dụng")

    cards = [
        ("01", "Nguồn video", "Nhận RTSP stream hoặc webcam để lấy frame realtime."),
        ("02", "Thu dataset", "Lưu ảnh khuôn mặt theo từng sinh viên từ camera hoặc thư mục ảnh."),
        ("03", "Train LBPH", "Sinh model nhận diện và file ánh xạ label sang MSSV, họ tên."),
        ("04", "Điểm danh", "Nhận diện khuôn mặt trong frame và ghi kết quả ra CSV theo ngày."),
    ]

    for idx, (num, title, body) in enumerate(cards):
        x = 0.72 + idx * 3.05
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(2.35), Inches(2.55), Inches(2.15))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.fill.transparency = 0.65
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        add_textbox(slide, x + 0.2, 2.6, 0.45, 0.2, num, size=13, color=ACCENT, bold=True, font=DISPLAY_FONT)
        add_textbox(slide, x + 0.2, 2.95, 2.1, 0.3, title, size=16, color=TEXT, bold=True, font=DISPLAY_FONT)
        add_textbox(slide, x + 0.2, 3.42, 2.05, 0.7, body, size=14, color=MUTED)

    add_textbox(
        slide,
        0.72,
        5.45,
        7.2,
        0.3,
        "Ứng dụng viết bằng Tkinter; mục tiêu là minh họa đầy đủ pipeline từ thu dữ liệu đến điểm danh.",
        size=14,
        color=MUTED,
    )
    add_rail(slide, 7, "USAGE")


def slide_summary():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_block(slide, "Takeaway", "Kết luận kỹ thuật")

    items = [
        ("01", "Haar Cascade là bộ phát hiện nhanh, phù hợp cho khuôn mặt chính diện và xử lý thời gian thực."),
        ("02", "Hiệu quả của thuật toán đến từ feature hình chữ nhật, integral image và cấu trúc cascade nhiều stage."),
        ("03", "Trong hệ thống điểm danh này, Haar Cascade và LBPH tách vai trò rõ ràng: phát hiện trước, nhận diện sau."),
    ]

    y = 2.45
    for num, text in items:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.35), Inches(y + 0.28), Inches(11.55), Inches(y + 0.28))
        line.line.color.rgb = LINE
        line.line.width = Pt(1)
        add_textbox(slide, 0.72, y, 0.45, 0.25, num, size=13, color=ACCENT, bold=True, font=DISPLAY_FONT)
        add_textbox(slide, 1.45, y, 8.9, 0.42, text, size=16, color=MUTED)
        y += 0.92

    add_textbox(
        slide,
        0.72,
        5.7,
        8.0,
        0.42,
        "Đây là lựa chọn phù hợp khi cần pipeline nhẹ, dễ triển khai và dễ giải thích trong bối cảnh bài tập hoặc demo học phần.",
        size=16,
        color=TEXT,
    )
    add_rail(slide, 8, "SUMMARY")


slide_intro()
slide_concept()
slide_features()
slide_training()
slide_cascade()
slide_pipeline()
slide_usage()
slide_summary()

output_path = Path(__file__).resolve().parent / "haar_cascade_attendance_deck_technical.pptx"
prs.save(output_path)
